#!/usr/bin/env python3
"""
Create PowerPoint presentation for Plant Segmentation project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define colors
TEAL_DARK = RGBColor(5, 68, 94)      # #05445E
TEAL_MED = RGBColor(8, 131, 149)     # #088395
TEAL_LIGHT = RGBColor(122, 178, 178) # #7AB2B2
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(235, 244, 246)       # #EBF4F6

def add_title_slide(title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = TEAL_DARK
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = TEAL_LIGHT
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(title, bullet_points, notes=""):
    """Add a content slide with bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = TEAL_DARK
    header.line.fill.background()

    # Title in header
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle sub-bullets (indicated by leading spaces)
        if point.startswith("  "):
            p.text = "    • " + point.strip()
            p.font.size = Pt(20)
            p.level = 1
        else:
            p.text = "• " + point
            p.font.size = Pt(24)
            p.level = 0

        p.font.color.rgb = TEAL_DARK
        p.space_after = Pt(12)

    # Add notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def add_two_column_slide(title, left_points, right_points, left_title="", right_title=""):
    """Add a two-column slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = TEAL_DARK
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column title
    if left_title:
        lt_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.5))
        tf = lt_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEAL_MED

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.8), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(20)
        p.font.color.rgb = TEAL_DARK
        p.space_after = Pt(8)

    # Right column title
    if right_title:
        rt_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.5), Inches(0.5))
        tf = rt_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEAL_MED

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(5.8), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(20)
        p.font.color.rgb = TEAL_DARK
        p.space_after = Pt(8)

    return slide

def add_screenshot_slide(title, caption="", notes=""):
    """Add a slide with placeholder for screenshot."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = TEAL_DARK
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Screenshot placeholder box
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.3), Inches(11.333), Inches(5.2)
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = PALE
    placeholder.line.color.rgb = TEAL_MED
    placeholder.line.width = Pt(2)

    # Placeholder text
    ph_text = slide.shapes.add_textbox(Inches(4), Inches(3.5), Inches(5.333), Inches(0.8))
    tf = ph_text.text_frame
    p = tf.paragraphs[0]
    p.text = "[ Insert Screenshot Here ]"
    p.font.size = Pt(24)
    p.font.color.rgb = TEAL_MED
    p.alignment = PP_ALIGN.CENTER

    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.6))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = TEAL_DARK
        p.alignment = PP_ALIGN.CENTER

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

# ═══════════════════════════════════════════════════════════════════════════
# CREATE SLIDES
# ═══════════════════════════════════════════════════════════════════════════

# 1. Title Slide
add_title_slide(
    "Plant Image Segmentation",
    "Interactive GUI for SAM2-based Leaf Segmentation & Phenotyping"
)

# 2. Overview
add_content_slide(
    "Project Overview",
    [
        "Desktop GUI application for plant phenotyping",
        "Uses SAM2 (Segment Anything Model 2) for automatic mask generation",
        "Interactive tools for mask editing and refinement",
        "Leaf shape completion for occluded leaves",
        "Export phenotype metrics (area, length, width, color stats)",
        "Batch processing for high-throughput analysis"
    ]
)

# 3. What is SAM?
add_content_slide(
    "What is SAM? (Segment Anything Model)",
    [
        "Foundation model for image segmentation by Meta AI (2023)",
        "Trained on 11 million images and 1.1 billion masks",
        "Can segment ANY object without task-specific training",
        "Promptable: points, boxes, or text can guide segmentation",
        "SAM2 (2024): Improved architecture with video support",
        "  Uses Hiera vision transformer backbone",
        "  Better quality masks with fewer parameters",
        "  Supports real-time interactive segmentation"
    ]
)

# 4. SAM Architecture
add_two_column_slide(
    "SAM2 Architecture",
    [
        "Image Encoder (Hiera ViT)",
        "  Processes image into embeddings",
        "  Hierarchical vision transformer",
        "  Pre-computed for efficiency",
        "",
        "Prompt Encoder",
        "  Encodes points, boxes, masks",
        "  Sparse (points) + dense (masks)",
        "  Positional encodings"
    ],
    [
        "Mask Decoder",
        "  Lightweight transformer",
        "  Generates mask predictions",
        "  Outputs multiple valid masks",
        "",
        "Memory Attention (SAM2)",
        "  Tracks objects across frames",
        "  Enables video segmentation",
        "  Temporal consistency"
    ],
    "Components", "Components (cont.)"
)

# 5. How SAM Works
add_content_slide(
    "How SAM Segmentation Works",
    [
        "Automatic Mask Generation Mode:",
        "  Grid of points sampled across image",
        "  Each point generates candidate masks",
        "  Non-maximum suppression removes duplicates",
        "  Filters by predicted IoU and stability score",
        "",
        "Key Parameters:",
        "  points_per_side: Density of point grid (default: 32)",
        "  pred_iou_thresh: Minimum IoU confidence (default: 0.88)",
        "  stability_score_thresh: Mask stability filter (default: 0.95)",
        "  min_mask_region_area: Remove tiny masks"
    ]
)

# 6. App Screenshot placeholder
add_screenshot_slide(
    "Plant Segmenter Application",
    "Main interface showing image preview, masks panel, and controls",
    "Take screenshot of the main app window showing the full interface"
)

# 7. Main Features
add_content_slide(
    "Application Features",
    [
        "Model Panel: Load SAM2 bundles or custom trained models",
        "Image Controls: Brightness, contrast, saturation, sharpness",
        "Automatic Segmentation: One-click mask generation",
        "Mask Management: Select, delete, combine, refine masks",
        "Mask Editing: Edit boundaries by dragging control points",
        "Leaf Completion: Complete partial/occluded leaf shapes",
        "Leaf Unfolding: Unfold folded leaves by rotating segments",
        "Batch Processing: Process entire folders automatically",
        "Export: Save masks and phenotype metrics to CSV"
    ]
)

# 8. Mask Editing Tools
add_two_column_slide(
    "Mask Editing Tools",
    [
        "🗑 Delete - Remove selected masks",
        "✖ Clear - Remove all masks",
        "🔗 Combine - Merge masks into one",
        "🔍 Refine - Re-segment within region",
        "✏ Edit - Edit mask boundary"
    ],
    [
        "Edit Boundary Features:",
        "Drag points to reshape",
        "Ctrl/Cmd+click to add points",
        "Shift+click to delete points",
        "Smooth button for curves",
        "Reset to original shape"
    ],
    "Toolbar Buttons", "Contour Editor"
)

# 9. Edit Mask Screenshot
add_screenshot_slide(
    "Edit Mask Boundary",
    "Contour editor with draggable control points for precise mask editing",
    "Take screenshot of the Edit Mask dialog showing control points"
)

# 10. Leaf Completion
add_content_slide(
    "Leaf Completion",
    [
        "Purpose: Complete partial or occluded leaf masks",
        "",
        "Geometric Shape Fitting:",
        "  Fits ellipse/oval to mask contour using cv2.fitEllipse()",
        "  Supports multiple leaf shapes: Ellipse, Ovate, Obovate, etc.",
        "  Interactive adjustments: size, position, rotation",
        "",
        "Edit Shape Feature:",
        "  Fine-tune fitted shape by dragging control points",
        "  Add/delete points for precise boundary control",
        "  Smooth function for natural curves",
        "  Preview shows original (blue) vs completed (green)"
    ]
)

# 11. Leaf Completion Screenshot
add_screenshot_slide(
    "Leaf Completion Interface",
    "Original mask preview and completed shape with geometric fitting",
    "Take screenshot of Leaf Completion tab showing before/after"
)

# 12. Edit Shape Screenshot
add_screenshot_slide(
    "Edit Completed Shape",
    "Fine-tune the fitted shape by dragging boundary points",
    "Take screenshot of the Edit Shape dialog for leaf completion"
)

# 13. Leaf Unfolding
add_content_slide(
    "Leaf Unfolding",
    [
        "Purpose: Unfold folded leaves for accurate measurements",
        "",
        "How it works:",
        "  Select 2+ masks representing parts of a folded leaf",
        "  Masks are combined in the preview",
        "  Select which mask to rotate",
        "  Adjust rotation angle (5° steps or quick 90°/180°)",
        "  Choose pivot point: auto-detected junction or mask center",
        "",
        "Result: Single unfolded mask for accurate phenotyping"
    ]
)

# 14. ML Shape Completion
add_content_slide(
    "ML-Based Shape Completion",
    [
        "Alternative to geometric fitting for complex cases",
        "",
        "Architecture: U-Net (encoder-decoder with skip connections)",
        "  Input: Partial binary mask (128×128)",
        "  Output: Completed binary mask",
        "",
        "Training:",
        "  Dataset: Complete leaf masks",
        "  Synthetic occlusion: 30-50% randomly removed",
        "  Loss: BCE + Dice loss combination",
        "  Augmentation: rotation, flip, scale",
        "",
        "Results: 95-98% IoU on synthetic test set"
    ]
)

# 15. Training Data Pipeline
add_two_column_slide(
    "Shape Completion Training",
    [
        "Data Preparation:",
        "Complete leaf masks collected",
        "Synthetic occlusion applied",
        "Random shapes overlay masks",
        "Leaf-shaped occlusions added",
        "30-65% mask removed"
    ],
    [
        "Training Process:",
        "U-Net learns mask→mask mapping",
        "No RGB needed, just binary masks",
        "BCE + Dice loss function",
        "15,000 training steps",
        "Model saved as .pth file"
    ],
    "Input Data", "Model Training"
)

# 16. Phenotype Metrics
add_content_slide(
    "Phenotype Metrics Export",
    [
        "Geometry Metrics:",
        "  Area (pixels and calibrated units)",
        "  Bounding box (x, y, width, height)",
        "  PCA-derived length and width",
        "  Perimeter and circularity",
        "",
        "Color Statistics:",
        "  RGB channel means, medians, std dev",
        "  HSV channel statistics",
        "  Useful for health/stress analysis",
        "",
        "Export: CSV format with all selected masks"
    ]
)

# 17. Technical Stack
add_two_column_slide(
    "Technical Implementation",
    [
        "Core Libraries:",
        "Python 3.11+",
        "PyTorch (model inference)",
        "OpenCV (image processing)",
        "NumPy (array operations)",
        "Tkinter (GUI framework)"
    ],
    [
        "Key Algorithms:",
        "SAM2 automatic mask generation",
        "cv2.fitEllipse() for shape fitting",
        "cv2.approxPolyDP() for contours",
        "PCA for orientation analysis",
        "Connected components for splits"
    ],
    "Dependencies", "Algorithms"
)

# 18. Workflow Summary
add_content_slide(
    "Typical Workflow",
    [
        "1. Load SAM2 model bundle",
        "2. Open image or browse folder",
        "3. Click 'Segment' for automatic mask generation",
        "4. Review and select relevant masks",
        "5. Edit masks if needed (combine, refine, edit boundary)",
        "6. Use Leaf Completion for occluded leaves",
        "7. Use Leaf Unfolding for folded leaves",
        "8. Export selected masks and metrics to CSV",
        "9. Repeat for batch processing"
    ]
)

# 19. Future Work
add_content_slide(
    "Future Improvements",
    [
        "Fine-tuning SAM2 on plant-specific datasets",
        "Integration with ImageJ/Fiji for advanced analysis",
        "3D reconstruction from multiple views",
        "Time-series growth tracking",
        "Disease detection and classification",
        "Cloud deployment for collaborative work",
        "Mobile app for field phenotyping"
    ]
)

# 20. Summary
add_title_slide(
    "Summary",
    "SAM2-powered segmentation • Interactive editing • Shape completion • Phenotype export"
)

# 21. Thank You
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Background
background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
)
background.fill.solid()
background.fill.fore_color.rgb = TEAL_DARK
background.line.fill.background()

# Thank you text
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Thank You!"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# GitHub link
link_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.8))
tf = link_box.text_frame
p = tf.paragraphs[0]
p.text = "Questions?"
p.font.size = Pt(28)
p.font.color.rgb = TEAL_LIGHT
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = "/Users/nirwantandukar/Documents/Github/Image-Segmentation/Plant_Segmentation_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
